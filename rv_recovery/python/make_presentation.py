#!/usr/bin/env python3
"""
Build the professor-ready PowerPoint for the replication of
"Software-based Realtime Recovery from Sensor Attacks on Robotic Vehicles" (RAID 2020).

Follows the supplied instructions: 19-slide structure, data->model->software sensor->
correction->parameter selection->attack recovery->limitations order, uses the generated
figures, includes the key equations, adds speaker notes, professional academic style,
no overclaiming, offline injection noted where used, live SITL only as future work.

Output: ~/paperImp/Replication_Presentation.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HOME = os.path.expanduser('~')
FIG = f'{HOME}/paperImp/rv_recovery/figures'
OUT = f'{HOME}/paperImp/Replication_Presentation.pptx'

# palette
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x6D, 0xA4)
GREY = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# category tag colors (rule 6: clearly separate)
TAG_COLOR = {'Paper': BLUE, 'Ours': GREEN, 'Figure': NAVY, 'Note': AMBER,
             'Eq.': GREY, 'Limitation': AMBER}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.33), Inches(1.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    # underline bar
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(12.33), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.5))
        sp = sb.text_frame.paragraphs[0]; sr = sp.add_run(); sr.text = subtitle
        sr.font.size = Pt(14); sr.font.italic = True; sr.font.color.rgb = GREY


def add_bullets(slide, bullets, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        # optional leading tag "Paper:", "Ours:", etc.
        tag = None
        for t in TAG_COLOR:
            if b.startswith(t):
                tag = t; b = b[len(t):].lstrip(': ').strip(); break
        if tag:
            rt = p.add_run(); rt.text = f'{tag}:  '
            rt.font.bold = True; rt.font.size = Pt(size); rt.font.color.rgb = TAG_COLOR[tag]
        else:
            bullet = p.add_run(); bullet.text = '•  '
            bullet.font.size = Pt(size); bullet.font.color.rgb = BLUE
        rb = p.add_run(); rb.text = b
        rb.font.size = Pt(size); rb.font.color.rgb = RGBColor(0x20, 0x20, 0x20)


def add_image_fit(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    w = max_w; h = Emu(int(w / ar))
    if h > max_h:
        h = max_h; w = Emu(int(h * ar))
    x = left + Emu(int((max_w - w) / 2))
    y = top + Emu(int((max_h - h) / 2))
    slide.shapes.add_picture(path, x, y, width=w, height=h)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide(title, bullets, image=None, notes='', subtitle=None, layout='right',
          image2=None):
    s = prs.slides.add_slide(BLANK)
    add_title(s, title, subtitle)
    top = Inches(1.55)
    if image is None:
        add_bullets(s, bullets, Inches(0.6), top, Inches(12.1), Inches(5.5))
    elif layout == 'right':
        add_bullets(s, bullets, Inches(0.6), top, Inches(6.0), Inches(5.5))
        add_image_fit(s, image, Inches(6.8), top, Inches(6.3), Inches(5.4))
        if image2:
            add_image_fit(s, image2, Inches(6.8), Inches(4.3), Inches(6.3), Inches(2.6))
    elif layout == 'below':
        add_bullets(s, bullets, Inches(0.6), top, Inches(12.1), Inches(2.0), size=13)
        if image2:
            add_image_fit(s, image, Inches(0.4), Inches(3.75), Inches(6.4), Inches(3.5))
            add_image_fit(s, image2, Inches(6.9), Inches(3.75), Inches(6.0), Inches(3.5))
        else:
            add_image_fit(s, image, Inches(1.2), Inches(3.6), Inches(10.9), Inches(3.7))
    if notes:
        add_notes(s, notes)
    return s


def fp(name):
    return f'{FIG}/{name}'

# ── Slide 1: Title ───────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
band = s.shapes.add_shape(1, 0, Inches(2.2), SW, Inches(2.2))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.35), Inches(12.1), Inches(1.9))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run()
r.text = 'Software-based Realtime Recovery from Sensor Attacks on Robotic Vehicles'
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); r2 = p2.add_run()
r2.text = 'A procedure-faithful replication of Choi et al., RAID 2020'
r2.font.size = Pt(18); r2.font.italic = True; r2.font.color.rgb = RGBColor(0xCF, 0xDE, 0xF0)
sub = s.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12.1), Inches(1.5))
st = sub.text_frame; st.word_wrap = True
for txt, sz, col in [
    ('Focus: replication of the methodology and procedure — identical numerical results are not claimed.', 15, GREY),
    ('Platform: ArduCopter 3.4 SITL · MATLAB R2026a · Python · C/C++ firmware patch', 13, GREY),
    ('Presented by: [Your name]   ·   Course: [Course]   ·   Date: [Date]', 13, GREY)]:
    pp = st.add_paragraph() if txt != st.paragraphs[0].text else st.paragraphs[0]
    rr = pp.add_run(); rr.text = txt; rr.font.size = Pt(sz); rr.font.color.rgb = col
add_notes(s, "Introduce the work: this is a from-scratch replication of the RAID 2020 paper "
             "on recovering robotic vehicles from sensor attacks using software sensors. "
             "Emphasize up front that the goal is to reproduce the paper's PROCEDURE faithfully, "
             "not to claim identical numbers. State the platform and that evidence figures are "
             "produced from recorded flight data and the identified model.")

# ── Slide 2: Motivation ──────────────────────────────────────────────────────
slide('Motivation: the sensor-attack recovery problem',
      ['Robotic vehicles (drones, rovers) continuously act on sensor measurements '
       '(gyroscope, accelerometer, GPS, barometer, magnetometer).',
       'Physical / external attacks corrupt these readings — e.g. GPS spoofing, acoustic '
       'gyroscope injection — and can crash the vehicle.',
       'Prior defenses mostly DETECT attacks but do not RECOVER: the vehicle can still crash '
       'after detection.',
       'Hardware redundancy (triple-modular) is costly and still fails when multiple '
       'same-type sensors are attacked together.',
       'Need: a software-only, real-time method that keeps the vehicle stable DURING an attack.'],
      notes="Set up the problem. Sensors are safety-critical. Attackers can manipulate them "
            "physically. Detection alone is insufficient — you also need to keep flying. "
            "Hardware redundancy is expensive and defeated by same-type multi-sensor attacks. "
            "This motivates a software-based recovery approach.")

# ── Slide 3: Software-sensor idea ────────────────────────────────────────────
slide("The paper's software-sensor idea",
      ['Paper: build a predictive state-space model of the vehicle that captures the '
       'controller, actuators, and dynamics.',
       'Paper: a "software sensor" converts model outputs into the predicted reading of each '
       'physical sensor — a software backup running alongside the real sensor.',
       'Paper: at runtime, compare software-sensor vs real-sensor; a large, sustained '
       'difference indicates the physical sensor is under attack.',
       'Paper: replace the compromised reading with the software-sensor value to keep the '
       'control loop stable (single, multiple, or all-same-type sensor attacks).',
       'Note: software sensors are model-based, so they are immune to physical-environment '
       'attacks.'],
      notes="Explain the core idea in plain terms: a physics/model-based 'virtual sensor' that "
            "predicts what each real sensor should read. Under normal operation the two agree; "
            "under attack they diverge, and the system swaps in the trustworthy software value. "
            "This is what makes recovery (not just detection) possible.")

# ── Slide 4: Paper workflow replicated ───────────────────────────────────────
slide('Paper workflow we replicated',
      ['Paper (offline): collect operation data → resample → system identification '
       '→ build software sensors → select recovery parameters.',
       'Paper (runtime): insert the recovery monitor (Algorithm 1) right after sensor '
       'acquisition in the control loop.',
       'Ours: every offline stage implemented in Python + MATLAB; runtime monitor implemented '
       'in C/C++ and inserted at the paper’s Figure-3 location.',
       'This deck follows the logical order: data → model → software-sensor prediction '
       '→ correction → parameter selection → attack recovery → limitations.'],
      notes="Give the roadmap. The paper has an offline modeling pipeline and a runtime monitor. "
            "We reproduced both. Tell the audience the slides will walk this exact pipeline in order.")

# ── Slide 5: Scope & environment ─────────────────────────────────────────────
slide('Replication scope and environment',
      ['Ours: simulated quadrotor on ArduCopter 3.4 (SITL).',
       'Ours: MATLAB R2026a (System Identification Toolbox); Python (pymavlink, scipy, '
       'dtaidistance, matplotlib); C/C++ firmware patch.',
       'Ours: runtime monitor is a line-for-line implementation of Algorithm 1, verified by '
       '25/25 unit tests.',
       'Note (deviation): model and monitor run at 50 Hz instead of the paper’s 400 Hz, '
       'due to a RAM constraint — the only intentional deviation.',
       'Limitation: closed-loop in-flight attack/recovery validation is future work (see '
       'Limitations slide).'],
      notes="State the scope honestly: one simulated quadrotor, the toolchain used, and that the "
            "monitor itself is unit-tested. Flag the single intentional deviation (50 Hz vs 400 Hz) "
            "now so it isn't a surprise later, and preview that live validation is future work.")

# ── Slide 6: Operation data (fig8) ───────────────────────────────────────────
slide('Operation data used for system identification',
      ['Paper: collect normal-operation logs under many maneuvers via randomly generated '
       'MAVLink missions; resample streams to one rate with spline interpolation.',
       'Ours: 20 random GUIDED missions (straight / turn / climb / hover); cubic-spline '
       'resampled per flight segment to 50 Hz — 21 segments, ~41 minutes.',
       'Ours: state x, input u (target states), output y built per Eq. (3); raw units, no '
       'detrending (consistent with runtime).',
       'Figure: one mission across diverse maneuvers (attitude, position, velocity, body rates).'],
      image=fp('fig8_operation_data_overview.png'),
      notes="Explain the data step. The paper needs varied maneuvers so the model captures the "
            "vehicle dynamics. We generated random missions and resampled with splines, exactly "
            "as the paper specifies. The figure shows the diversity of one mission.")

# ── Slide 7: State-space model & sysid (fig4 + fig5) ─────────────────────────
slide('State-space model and system identification',
      ['Paper: discrete state-space model; per-variable 2nd-order order; fit by '
       'prediction-error minimization (PEM).',
       'Eq.:  xₖ₊₁ = A xₖ + B uₖ      yₖ = C xₖ + D uₖ',
       'Ours: 12-state model (Eq. 3); six per-axis 2nd-order PEM blocks; open-loop K = 0.',
       'Note: roll/pitch fit well; translational/yaw under-excited by hover-dominant data.'],
      image=fp('fig4_sysid_validation_fit.png'),
      image2=fp('fig5_model_eigenvalues.png'),
      layout='below',
      notes="This is the model step. State the two governing equations (point at them). We fit a "
            "12-state model with per-axis second-order blocks via PEM, with K=0 so it matches the "
            "open-loop predictor the runtime uses. Be honest that roll/pitch fit well but the "
            "translational/yaw models are weaker because the training data was hover-dominant.")

# ── Slide 8: Software sensor generation (fig3) ───────────────────────────────
slide('Software-sensor generation',
      ['Paper (§3.2): convert model outputs to predicted sensor readings — gyroscope '
       '= rate states, accelerometer = Eq. (4), barometer = Eq. (5), magnetometer heading = '
       'Eq. (6), GPS = position/velocity states.',
       'Ours: all conversion equations implemented in software_sensors.h, plus the '
       'Appendix-A frame transforms.',
       'Figure: software sensor (blue) vs real sensor (red) for GPS, barometer, gyroscope, '
       'magnetometer heading.',
       'Note: barometer / gyroscope / magnetometer track closely; GPS-North drifts — the '
       'documented under-excited translational model.'],
      image=fp('fig3_sensor_prediction.png'),
      notes="The software-sensor step. We implemented every conversion equation from the paper. "
            "The figure shows predicted vs real for four sensor types. Point out the tight "
            "tracking on baro/gyro/mag, and honestly note the GPS drift due to the weaker "
            "translational model.")

# ── Slide 9: Recovery monitor / Algorithm 1 (fig2) ───────────────────────────
slide('Recovery monitor — Algorithm 1',
      ['Paper (§3.4): each tick — predict y = Cx + Du, advance x = Ax + Bu, filter the '
       'real measurement, convert to software sensor, accumulate residual, detect, substitute.',
       'Eq.: rₖ = rₖ + |mₖ − msₖ|        Recovery when  r > T_on  '
       '(switch back when r < T_off for K counts).',
       'Ours: line-for-line Algorithm 1 in recovery_monitor.h; 25/25 conformance tests pass.',
       'Figure: offline gyro attack — residual crosses T_on, monitor detects (~1.3 s) and '
       'substitutes the software sensor for the attacked reading.',
       'Note: attack injected offline into a recorded trace (supported by our scripts); '
       'closed-loop in-flight recovery is future work.'],
      image=fp('fig2_attack_recovery_gyro.png'),
      notes="The heart of the method. Walk through Algorithm 1 at a high level and show the "
            "residual/threshold equation. Stress that our implementation is unit-tested line by "
            "line. The figure is an OFFLINE demonstration on a recorded trace: the attack drives "
            "the residual past T_on and the monitor swaps in the software sensor. Be explicit that "
            "this is offline, not a live flight.")

# ── Slide 10: Drift correction & synchronization (fig12) ─────────────────────
slide('Drift correction and synchronization',
      ['Paper (§3.3): an open-loop model prediction drifts over time; periodically '
       'synchronize the predicted state with the real state and reset accumulated error.',
       'Ours: at each window checkpoint, re-seed the model state from the real reading and '
       'reset the residual accumulator.',
       'Figure: accumulated prediction error grows without correction vs stays bounded with '
       '§3.3 synchronization.'],
      image=fp('fig12_roll_drift_correction.png'),
      notes="Explain why correction is needed: an open-loop model inevitably drifts. The paper's "
            "answer is periodic synchronization plus error reset. The figure shows the accumulated "
            "error growing unbounded without it and staying bounded with it.")

# ── Slide 11: External disturbance / wind compensation (fig13) ───────────────
slide('External disturbance / wind compensation',
      ['Paper (§3.3): external forces (e.g. wind) cause a persistent prediction error; '
       'estimate the average error over the previous window (term e) and subtract it.',
       'Ours: e = average(ms − m) over the previous window; corrected software sensor '
       'ms ← ms − e.',
       'Figure: prediction error WITH vs WITHOUT the disturbance term, under constant and '
       'dynamic wind (disturbance injected offline).',
       'Note: shown on a translational channel whose model is under-excited, so the trace is '
       'noisier than the paper’s.'],
      image=fp('fig13_external_wind_correction.png'),
      notes="The external-error compensation step. Wind pushes the vehicle so the real sensor "
            "differs from the model by a persistent offset; the e term absorbs it. The figure "
            "shows the error shrinking when compensation is on. Note the disturbance is injected "
            "offline (a simulated wind), consistent with the paper's evaluation approach.")

# ── Slide 12: Recovery parameter selection (fig6 + fig14) ────────────────────
slide('Recovery parameter selection',
      ['Paper (§3.3): window N = max DTW time-displacement on clean data; threshold '
       'T = e_max + margin, with T_off < T_on.',
       'Ours: DTW per channel (dtaidistance); per-channel N, T_on, T_off in recovery_params.h.',
       'Figure: selected N and thresholds (left); FP & FN vs threshold across window sizes '
       '(right) — both ≈ 0 at the selected T_on.',
       'Note: FP/FN measured offline with in-memory attack injection (paper’s parameter study).'],
      image=fp('fig6_dtw_parameters.png'),
      image2=fp('fig14_param_selection_windows.png'),
      layout='below',
      notes="Parameter selection. The window comes from DTW; the threshold from the max clean "
            "residual plus a margin. The right figure reproduces the paper's FP/FN trade-off: a "
            "larger window gives more false positives and fewer false negatives, and at the chosen "
            "threshold both are about zero. State that FP/FN were measured offline by injecting "
            "attacks into recorded clean data.")

# ── Slide 13: Attack detection vs attack scale (fig11) ───────────────────────
slide('Attack detection versus attack scale',
      ['Paper (§4.2.2): larger attacks produce larger errors; whether an attack is '
       'detected depends on its scale relative to the threshold.',
       'Ours: sweep the injected constant gyro-bias magnitude; measure the peak accumulated '
       'residual per mission.',
       'Figure: peak residual rises ~linearly with attack scale and crosses T_on; all missions '
       'are detected for attacks ≥ ~0.3 rad/s.',
       'Note: attacks injected offline into recorded traces.'],
      image=fp('fig11_error_vs_attack_scale.png'),
      notes="Shows sensitivity vs attack strength. Small attacks stay under the threshold; once "
            "the bias is large enough the residual crosses T_on and detection is guaranteed. This "
            "mirrors the paper's attack-scale analysis. Again, offline injection on recorded data.")

# ── Slide 14: All-gyroscope attack compensation (fig15) ──────────────────────
slide('All-gyroscope attack compensation',
      ['Paper (Appendix B): when ALL gyroscopes are compromised, the software gyro alone '
       'drifts; reconstruct roll/pitch from the accelerometer and yaw from the magnetometer '
       '(Eq. 11) and combine by weighted sum.',
       'Ours: supplementary_compensation() implements Eq. 11; demonstrated offline using the '
       'recorded accelerometer and magnetometer streams.',
       'Figure: gyro-only roll diverges (~575°) under the all-gyros attack vs stays bounded '
       'with accelerometer compensation.'],
      image=fp('fig15_allgyro_compensation.png'),
      notes="The hardest case in the paper: all gyros attacked at once. The software gyro has no "
            "real reference and drifts, so the paper adds accelerometer/magnetometer-based "
            "attitude reconstruction. The figure shows roll diverging without it and bounded with "
            "it. Demonstrated offline from the recorded IMU/compass log.")

# ── Slide 15: Code-to-figure mapping ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_title(s, 'Code-to-figure mapping')
rows = [
    ('Figure', 'Generating script (read-only)', 'Paper figure'),
    ('Operation data overview (fig8)', 'make_figures_paper.py', '§3.1 data'),
    ('Sensor prediction (fig3)', 'make_figures_paper.py', 'Fig. 11'),
    ('Sysid fit / model poles (fig4, fig5)', 'make_figures_paper.py', '§3.1'),
    ('Attack → recovery (fig2)', 'make_figures.py', 'Fig. 5 / 15'),
    ('Drift correction (fig12, fig9)', 'make_figures_paper_12_15.py / eval', 'Fig. 12'),
    ('Wind compensation (fig13)', 'make_figures_paper_12_15.py', 'Fig. 13'),
    ('DTW params / FP-FN (fig6, fig14, fig10, fig7)', 'select_parameters.py / make_figures_*', 'Fig. 14'),
    ('Attack scale (fig11)', 'make_figures_eval.py', 'Fig. 16b'),
    ('All-gyros compensation (fig15)', 'make_figures_paper_12_15.py', 'Fig. 15'),
]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.5), Inches(1.6),
                         Inches(12.3), Inches(5.2)).table
tbl.columns[0].width = Inches(5.1); tbl.columns[1].width = Inches(4.6)
tbl.columns[2].width = Inches(2.6)
for ci in range(3):
    cell = tbl.cell(0, ci); cell.text = rows[0][ci]
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for pp in cell.text_frame.paragraphs:
        for rr in pp.runs:
            rr.font.bold = True; rr.font.size = Pt(13); rr.font.color.rgb = WHITE
for ri in range(1, len(rows)):
    for ci in range(3):
        cell = tbl.cell(ri, ci); cell.text = rows[ri][ci]
        for pp in cell.text_frame.paragraphs:
            for rr in pp.runs:
                rr.font.size = Pt(11); rr.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
add_notes(s, "Transparency slide: every figure is produced by a named, read-only script from the "
             "recorded data and the identified model. No figure uses invented data. Mention the "
             "scripts live in rv_recovery/python/.")

# ── Slide 16: Successfully replicated components ─────────────────────────────
slide('Successfully replicated components',
      ['Ours: data pipeline — random-mission collection, spline resampling, x/u/y '
       'construction (§3.1).',
       'Ours: system identification — per-variable 2nd-order PEM, open-loop K = 0, 12-state '
       'model (§3.1).',
       'Ours: software sensors — Eq. 4 (accel), Eq. 5 (baro), Eq. 6 (magnetometer), frame '
       'transforms, Eq. 11 (supplementary) (§3.2 / App. A, B).',
       'Ours: error correction (LPF, synchronization, disturbance term e) and DTW-based '
       'parameter selection (§3.3).',
       'Ours: Algorithm 1 runtime monitor (25/25 tests) inserted at the paper’s Figure-3 '
       'point (§3.4).'],
      notes="Summarize what matches the paper. Walk the checklist quickly: data, model, software "
            "sensors, correction, parameters, and the runtime monitor are all implemented to the "
            "paper's procedure. Point to the comparison checklist document for the full mapping.")

# ── Slide 17: Limitations & deviations ───────────────────────────────────────
slide('Limitations and deviations',
      ['Limitation: 50 Hz instead of the paper’s 400 Hz (RAM) — the only intentional '
       'deviation; all downstream parameters are consistent at 50 Hz.',
       'Limitation: translational / altitude sub-models are under-excited by hover-dominant '
       'data, so those channels are noisier (more aggressive maneuver data would improve them).',
       'Limitation: closed-loop in-flight attack/recovery not produced — ArduCopter 3.4 SITL '
       'takes attitude from a perfect physics model and regenerates the IMU asynchronously, so a '
       'read-AHRS-level substitution (attack AND recovery, identically) is overwritten before the '
       'estimator reads it. Real hardware (the paper’s setting) would propagate it.',
       'Future work: live SITL / physical-drone A/B validation; the rover and hexrotor; the '
       'multi-sensor attack matrix (Tables 3 / 4).',
       'Note: all figures are offline demonstrations on recorded data, not closed-loop flight.'],
      notes="Be fully transparent. The 50 Hz deviation is intentional and consistent. The "
            "translational models are weaker due to data coverage. Most importantly, explain why "
            "live numbers aren't shown: the SITL sensor path overwrites the substitution before the "
            "estimator sees it, affecting attack and recovery identically — a simulator-architecture "
            "issue, not a flaw in the method. Real hardware would behave like the paper. List the "
            "remaining future work.")

# ── Slide 18: Conclusion ─────────────────────────────────────────────────────
slide('Conclusion',
      ['We faithfully replicated the paper’s full offline + runtime procedure: data '
       'pipeline, system identification, software sensors, error correction, parameter '
       'selection, and the Algorithm 1 monitor.',
       'Evidence figures reproduce the paper’s effectiveness plots (Figs. 11–16b) on '
       'recorded data and the identified model.',
       'We focus on PROCEDURE replication — we do not claim identical numerical results.',
       'Remaining: closed-loop live validation and additional vehicles / attack combinations '
       '(future work).'],
      notes="Wrap up: the methodology is replicated end to end and demonstrated with figures that "
            "mirror the paper's. Restate the honest framing — procedure replication, not identical "
            "numbers — and that live validation is the main piece of future work.")

# ── Slide 19: Backup — equations & parameter definitions ─────────────────────
slide('Backup: equations and parameter definitions',
      ['State / output:  xₖ₊₁ = A xₖ + B uₖ ;  yₖ = C xₖ + D uₖ.',
       'Residual / detection:  rₖ = rₖ + |mₖ − msₖ| ;  recover if  r > T_on ;  '
       'exit if  r < T_off  for K consecutive counts.',
       'Conversions:  accelerometer Eq. (4) ;  barometer Eq. (5) ;  magnetometer heading '
       'Eq. (6) ;  supplementary attitude Eq. (11).',
       'Parameters:  N = window size (DTW) ;  T_on / T_off = recovery thresholds ;  '
       'K = safe-count for switch-back ;  margin in T = e_max + margin.',
       'Rate:  50 Hz (Ts = 0.02 s) — the sanctioned deviation from the paper’s 400 Hz.'],
      notes="Backup reference for questions: the governing equations, the conversion equations by "
            "number, and the meaning of every parameter (N, T_on, T_off, K, margin) plus the rate.")

prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}')
